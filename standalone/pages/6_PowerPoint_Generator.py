"""
Credentials Template Filler - Chat Agent

Natural language interface to fill PowerPoint templates with credentials.
Just ask: "Fill in the template with Bud van der Schier's most recent experience"
"""

import streamlit as st
import pandas as pd
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
from io import BytesIO
import requests
import json
from typing import Dict, List, Optional
import re

# ==== Workbench API Configuration ====
API_URL = "https://api.workbench.kpmg/genai/azure/inference/chat/completions?api-version=2024-12-01-preview"

DEFAULT_SUBSCRIPTION_KEY = "b82fef87872349b981d5c0d58afb55c1"
DEFAULT_CHARGE_CODE = "1"
DEFAULT_DEPLOYMENT = "gpt-4o-2024-08-06-dzs-we"

def get_api_config():
    """Get API configuration"""
    if "api_config" not in st.session_state:
        try:
            secrets = st.secrets.get("workbench", {})
            subscription_key = (
                secrets.get("subscription_key") or 
                os.environ.get("WORKBENCH_SUBSCRIPTION_KEY") or 
                DEFAULT_SUBSCRIPTION_KEY
            )
            charge_code = (
                secrets.get("charge_code") or 
                os.environ.get("WORKBENCH_CHARGE_CODE") or 
                DEFAULT_CHARGE_CODE
            )
            deployment = (
                secrets.get("deployment") or 
                os.environ.get("WORKBENCH_DEPLOYMENT") or 
                DEFAULT_DEPLOYMENT
            )
        except:
            subscription_key = os.environ.get("WORKBENCH_SUBSCRIPTION_KEY") or DEFAULT_SUBSCRIPTION_KEY
            charge_code = os.environ.get("WORKBENCH_CHARGE_CODE") or DEFAULT_CHARGE_CODE
            deployment = os.environ.get("WORKBENCH_DEPLOYMENT") or DEFAULT_DEPLOYMENT
        
        st.session_state.api_config = {
            "subscription_key": subscription_key,
            "charge_code": charge_code,
            "deployment": deployment
        }
    
    return st.session_state.api_config

def get_api_headers():
    """Get API headers"""
    config = get_api_config()
    return {
        'Ocp-Apim-Subscription-Key': config["subscription_key"],
        'Cache-Control': 'no-cache',
        'x-kpmg-charge-code': config["charge_code"],
        'azureml-model-deployment': config["deployment"]
    }

# Page configuration
st.set_page_config(
    page_title="Credentials Chat Agent",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded"
)

st.title("🤖 Credentials Template Filler - Chat Agent")
st.markdown("Just tell me what you want! Example: *'Fill in the template with Bud van der Schier's most recent experience'*")

@st.cache_data
def load_credentials_data():
    """Load credentials database"""
    # Try multiple paths for the Excel file
    possible_paths = [
        'credentials_full.xlsx',  # Same directory as app.py (when run from standalone/)
        os.path.join(os.path.dirname(__file__), '..', 'credentials_full.xlsx'),  # Parent directory
        os.path.join(os.path.dirname(os.path.dirname(__file__)), 'credentials_full.xlsx'),  # Two levels up
        os.path.join(os.path.dirname(__file__), 'credentials_full.xlsx')  # Pages directory
    ]
    
    for excel_path in possible_paths:
        try:
            if os.path.exists(excel_path):
                df = pd.read_excel(excel_path)
                return df
        except Exception as e:
            continue
    
    # If all attempts fail, show detailed error
    st.error(f"Error loading data: [Errno 2] No such file or directory: 'C:\\Users\\mmartini1\\KPMG_Credentials_System\\pages\\credentials_full.xlsx'")
    st.info("Please ensure 'credentials_full.xlsx' exists in the credentials folder")
    return None

def call_llm(system_prompt: str, user_message: str) -> Optional[str]:
    """Call the LLM API"""
    try:
        config = get_api_config()
        headers = get_api_headers()
        
        body = {
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_message}
            ],
            "max_completion_tokens": 1000,
            "temperature": 0.3
        }
        
        response = requests.post(API_URL, headers=headers, json=body, timeout=30)
        
        if response.status_code == 200:
            result = response.json()
            content = result['choices'][0]['message']['content']
            return content
        else:
            st.warning(f"API Error: {response.status_code}")
            return None
            
    except Exception as e:
        st.warning(f"Error calling LLM: {str(e)}")
        return None

def parse_user_request(user_request: str, available_people: List[str]) -> Optional[Dict]:
    """Parse natural language request to extract intent"""
    
    # First try simple keyword matching for common patterns
    user_lower = user_request.lower()
    
    # Look for person names in the request
    found_person = None
    for person in available_people:
        person_lower = person.lower()
        # Check if full name or last name is in the request
        if person_lower in user_lower:
            found_person = person
            break
        # Check last name only
        last_name = person.split()[-1].lower()
        if len(last_name) > 3 and last_name in user_lower:
            found_person = person
            break
    
    # If we found someone with simple matching, return immediately
    if found_person:
        return {
            "action": "generate_credentials",
            "person_name": found_person,
            "num_projects": 5,
            "preferences": "most recent",
            "found_match": True,
            "explanation": f"Generate credentials for {found_person}"
        }
    
    # If simple matching failed, try LLM
    people_list = ", ".join(available_people[:50])
    
    system_prompt = f"""You are an assistant that parses user requests for generating credential presentations.

Available people in the database:
{people_list}

The user will ask to fill in a template or generate credentials for someone. Extract:
1. The person's name (match against available people, use VERY FLEXIBLE fuzzy matching)
2. Number of projects (default: 5 if not specified)
3. Any preferences (most recent, specific industry, etc.)

IMPORTANT: Be VERY flexible with name matching. Match partial names, last names only, or similar spellings.

Return ONLY a JSON object:
{{
    "action": "generate_credentials",
    "person_name": "exact name from available people",
    "num_projects": 5,
    "preferences": "most recent / specific criteria",
    "found_match": true/false,
    "explanation": "what you understood"
}}

If you find a reasonable match, set found_match to true. Only set it to false if there's absolutely no match.
"""

    response = call_llm(system_prompt, user_request)
    
    if response:
        try:
            # Clean JSON from response
            if response.startswith("```json"):
                response = response[7:]
            if response.startswith("```"):
                response = response[3:]
            if response.endswith("```"):
                response = response[:-3]
            response = response.strip()
            
            return json.loads(response)
        except json.JSONDecodeError as e:
            st.error(f"Error parsing response: {str(e)}")
            return None
    
    return None

def find_person_projects(df: pd.DataFrame, person_name: str) -> pd.DataFrame:
    """Find all projects for a specific person"""
    partner_col = None
    manager_col = None
    
    for col in df.columns:
        col_lower = col.lower()
        if 'partner' in col_lower and 'engagement' in col_lower:
            partner_col = col
        elif 'manager' in col_lower and 'engagement' in col_lower:
            manager_col = col
    
    mask = pd.Series([False] * len(df))
    
    if partner_col:
        mask |= df[partner_col].astype(str).str.lower().str.contains(
            person_name.lower(), na=False
        )
    
    if manager_col:
        mask |= df[manager_col].astype(str).str.lower().str.contains(
            person_name.lower(), na=False
        )
    
    return df[mask].copy()

def select_projects_with_llm(projects_df: pd.DataFrame, num_projects: int, preferences: str = "most recent") -> Dict:
    """Use LLM to select the best projects"""
    
    project_summaries = []
    for idx, row in projects_df.iterrows():
        summary = {
            "id": idx,
            "project_name": row.get("Project name / engagement title", "N/A"),
            "client": row.get("Client name", "N/A"),
            "year": row.get("Year of completion", "N/A"),
            "sector": row.get("Sector", "N/A"),
            "industry": row.get("Industry", "N/A"),
            "service": row.get("Service offering / proposition", "N/A"),
            "description": str(row.get("Credential description including proposition (tombstone appropriate)", "N/A"))[:200]
        }
        project_summaries.append(summary)
    
    system_prompt = f"""You are an expert at selecting the most relevant and impressive project credentials.

Given {len(projects_df)} projects, select the TOP {num_projects} projects that are:
- {preferences}
- Most diverse (different industries/sectors)
- Most impactful (impressive descriptions)
- Good variety of services

Return ONLY a JSON object:
{{
    "selected_project_ids": [list of project IDs],
    "reasoning": "brief explanation"
}}

Projects:
{json.dumps(project_summaries, indent=2)}
"""

    response = call_llm(system_prompt, f"Select the top {num_projects} projects based on criteria: {preferences}")
    
    if response:
        try:
            if response.startswith("```json"):
                response = response[7:]
            if response.startswith("```"):
                response = response[3:]
            if response.endswith("```"):
                response = response[:-3]
            response = response.strip()
            
            return json.loads(response)
        except json.JSONDecodeError:
            return None
    
    return None

def apply_filters_with_llm(projects_df: pd.DataFrame, filter_request: str) -> pd.DataFrame:
    """Use LLM to interpret and apply filters to the project database"""
    
    # Get column information for the LLM
    columns_info = {}
    for col in projects_df.columns:
        col_lower = col.lower()
        if any(keyword in col_lower for keyword in ['year', 'sector', 'industry', 'service', 'client', 'project']):
            unique_vals = projects_df[col].dropna().unique()
            if len(unique_vals) < 50:  # Only include if manageable list
                columns_info[col] = list(unique_vals[:20])  # First 20 examples
    
    system_prompt = f"""You are a data filtering assistant. The user wants to filter a project database.

Available columns and sample values:
{json.dumps(columns_info, indent=2, default=str)}

User's filter request: "{filter_request}"

Analyze the request and return filtering criteria as JSON:
{{
    "filters": [
        {{
            "column": "exact column name",
            "operation": "equals/contains/greater_than/less_than/between/in_list",
            "value": "filter value or [list] or [min, max]",
            "case_sensitive": false
        }}
    ],
    "reasoning": "explanation of filters applied"
}}

Examples:
- "projects from 2024-2025" → year between [2024, 2025]
- "only financial services" → sector/industry contains "financial"
- "exclude audit projects" → service line NOT contains "audit"

Return ONLY valid JSON."""

    response = call_llm(system_prompt, filter_request)
    
    if not response:
        return projects_df
    
    try:
        # Clean JSON
        if response.startswith("```json"):
            response = response[7:]
        if response.startswith("```"):
            response = response[3:]
        if response.endswith("```"):
            response = response[:-3]
        response = response.strip()
        
        filter_config = json.loads(response)
        
        if 'filters' not in filter_config or not filter_config['filters']:
            return projects_df
        
        filtered_df = projects_df.copy()
        
        for filter_rule in filter_config['filters']:
            column = filter_rule.get('column')
            operation = filter_rule.get('operation')
            value = filter_rule.get('value')
            case_sensitive = filter_rule.get('case_sensitive', False)
            
            if column not in filtered_df.columns:
                continue
            
            if operation == 'equals':
                if case_sensitive:
                    filtered_df = filtered_df[filtered_df[column] == value]
                else:
                    filtered_df = filtered_df[filtered_df[column].astype(str).str.lower() == str(value).lower()]
            
            elif operation == 'contains':
                if case_sensitive:
                    filtered_df = filtered_df[filtered_df[column].astype(str).str.contains(str(value), na=False)]
                else:
                    filtered_df = filtered_df[filtered_df[column].astype(str).str.lower().str.contains(str(value).lower(), na=False)]
            
            elif operation == 'greater_than':
                filtered_df = filtered_df[pd.to_numeric(filtered_df[column], errors='coerce') > float(value)]
            
            elif operation == 'less_than':
                filtered_df = filtered_df[pd.to_numeric(filtered_df[column], errors='coerce') < float(value)]
            
            elif operation == 'between' and isinstance(value, list) and len(value) == 2:
                col_numeric = pd.to_numeric(filtered_df[column], errors='coerce')
                filtered_df = filtered_df[(col_numeric >= float(value[0])) & (col_numeric <= float(value[1]))]
            
            elif operation == 'in_list' and isinstance(value, list):
                if case_sensitive:
                    filtered_df = filtered_df[filtered_df[column].isin(value)]
                else:
                    filtered_df = filtered_df[filtered_df[column].astype(str).str.lower().isin([str(v).lower() for v in value])]
        
        st.info(f"**Filters applied:** {filter_config.get('reasoning', 'Custom filters')}\n\n**Results:** {len(filtered_df)} projects (from {len(projects_df)} total)")
        
        return filtered_df
        
    except (json.JSONDecodeError, Exception) as e:
        st.warning(f"Could not parse filters: {str(e)}. Returning all projects.")
        return projects_df

def generate_project_description_with_llm(project_info: Dict) -> str:
    """Use LLM to generate a concise, professional project description"""
    
    system_prompt = """You are a professional consultant writing credential descriptions for PowerPoint presentations.

Given project information, create a CONCISE, PROFESSIONAL one-line description (max 100 characters) that highlights the key value delivered.

Focus on:
- The main deliverable or outcome
- Business impact
- Technical expertise demonstrated

Use professional consulting language. Be specific and impressive.

Return ONLY the description text, nothing else."""

    project_summary = f"""
Project: {project_info.get('project_name', 'N/A')}
Client: {project_info.get('client', 'N/A')}
Industry/Sector: {project_info.get('industry', 'N/A')} / {project_info.get('sector', 'N/A')}
Service: {project_info.get('service', 'N/A')}
Original Description: {project_info.get('description', 'N/A')[:300]}
"""

    response = call_llm(system_prompt, project_summary)
    
    if response:
        # Clean up response
        desc = response.strip().strip('"').strip("'")
        # Ensure it's not too long
        if len(desc) > 120:
            desc = desc[:117] + "..."
        return desc
    
    # Fallback to original description if LLM fails
    original = str(project_info.get('description', ''))
    if len(original) > 100:
        return original[:97] + "..."
    return original

def select_best_projects(projects_df: pd.DataFrame, num_projects: int, preferences: str = "most recent") -> pd.DataFrame:
    """Select the best projects"""
    
    # First, sort by year to ensure we prioritize recent projects
    year_col = None
    for col in projects_df.columns:
        if 'year' in col.lower() and 'completion' in col.lower():
            year_col = col
            break
    
    if year_col:
        projects_sorted = projects_df.sort_values(by=year_col, ascending=False, na_position='last')
    else:
        projects_sorted = projects_df
    
    # Take top 20 most recent projects as candidates (or all if fewer than 20)
    candidate_pool = projects_sorted.head(min(20, len(projects_sorted)))
    
    # If requesting fewer than available candidates, let LLM choose the most diverse/relevant ones
    if num_projects < len(candidate_pool):
        llm_result = select_projects_with_llm(candidate_pool, num_projects, preferences)
        
        if llm_result and 'selected_project_ids' in llm_result:
            selected_ids = llm_result['selected_project_ids']
            return candidate_pool.loc[selected_ids], llm_result.get('reasoning', '')
    
    # Fallback: return most recent projects directly
    return projects_sorted.head(num_projects), "Selected most recent projects"

def create_filled_presentation(template_path: str, projects_df: pd.DataFrame, person_name: str, include_details: bool = True) -> BytesIO:
    """Create a filled PowerPoint presentation"""
    
    prs = Presentation(template_path)
    
    # Find the template slide (first slide)
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        
        # Find the "Relevant projects" text shape
        relevant_projects_shape = None
        for shape in slide.shapes:
            if hasattr(shape, "text") and "Relevant projects" in shape.text:
                relevant_projects_shape = shape
                break
        
        if relevant_projects_shape and hasattr(relevant_projects_shape, "text_frame"):
            tf = relevant_projects_shape.text_frame
            
            # Find where "Relevant projects" ends and replace only the content after it
            # Keep all existing paragraphs until we find "Relevant projects"
            relevant_projects_index = -1
            for i, paragraph in enumerate(tf.paragraphs):
                if "Relevant projects" in paragraph.text:
                    relevant_projects_index = i
                    break
            
            if relevant_projects_index >= 0:
                # Remove only the paragraphs after "Relevant projects" header
                # We need to work backwards to avoid index issues
                paragraphs_to_remove = len(tf.paragraphs) - relevant_projects_index - 1
                for _ in range(paragraphs_to_remove):
                    # Remove the last paragraph
                    if len(tf.paragraphs) > relevant_projects_index + 1:
                        p = tf.paragraphs[-1]
                        p._element.getparent().remove(p._element)
                
                # Now add the new project bullets - all info on one line per project
                for idx, row in projects_df.iterrows():
                    project_name = row.get("Project name / engagement title", "Project")
                    client_name = row.get("Client name", "")
                    year = row.get("Year of completion", "")
                    sector = row.get("Sector", "")
                    industry = row.get("Industry", "")
                    service = row.get("Service offering / proposition", "")
                    description = row.get("Credential description including proposition (tombstone appropriate)", "")
                    
                    # Build single-line project bullet - project title with year at the end
                    p = tf.add_paragraph()
                    
                    # Format: "Project title (Year)" - simpler, cleaner format
                    text_parts = [project_name]
                    
                    if year and year != "N/A" and str(year).strip():
                        text_parts.append(f" ({year})")
                    
                    # Only add description if details requested - use LLM for dynamic generation
                    # Don't add industry/sector to keep it clean like the example
                    if include_details:
                        project_info = {
                            'project_name': project_name,
                            'client': client_name,
                            'industry': industry,
                            'sector': sector,
                            'service': service,
                            'description': description
                        }
                        generated_desc = generate_project_description_with_llm(project_info)
                        if generated_desc and len(generated_desc.strip()) > 5:
                            text_parts.append(f" - {generated_desc}")
                    
                    p.text = "".join(text_parts)
                    p.level = 1  # All as main bullets, no sub-bullets
            else:
                # If we can't find the exact location, just append at the end
                for idx, row in projects_df.iterrows():
                    project_name = row.get("Project name / engagement title", "Project")
                    client_name = row.get("Client name", "")
                    year = row.get("Year of completion", "")
                    sector = row.get("Sector", "")
                    industry = row.get("Industry", "")
                    service = row.get("Service offering / proposition", "")
                    description = row.get("Credential description including proposition (tombstone appropriate)", "")
                    
                    # Build single-line project bullet - project title with year at the end
                    p = tf.add_paragraph()
                    
                    # Format: "Project title (Year)" - simpler, cleaner format
                    text_parts = [project_name]
                    
                    if year and year != "N/A" and str(year).strip():
                        text_parts.append(f" ({year})")
                    
                    # Only add description if details requested - use LLM for dynamic generation
                    # Don't add industry/sector to keep it clean like the example
                    if include_details:
                        project_info = {
                            'project_name': project_name,
                            'client': client_name,
                            'industry': industry,
                            'sector': sector,
                            'service': service,
                            'description': description
                        }
                        generated_desc = generate_project_description_with_llm(project_info)
                        if generated_desc and len(generated_desc.strip()) > 5:
                            text_parts.append(f" - {generated_desc}")
                    
                    p.text = "".join(text_parts)
                    p.level = 1  # All as main bullets, no sub-bullets
    
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    
    return pptx_buffer

# Main application
df = load_credentials_data()

if df is not None and not df.empty:
    
    # Get available people
    partners = []
    managers = []
    
    for col in df.columns:
        col_lower = col.lower()
        if 'partner' in col_lower and 'engagement' in col_lower:
            partners = sorted([str(x) for x in df[col].dropna().unique()])
        elif 'manager' in col_lower and 'engagement' in col_lower:
            managers = sorted([str(x) for x in df[col].dropna().unique()])
    
    all_people = sorted(set(partners + managers))
    
    # Initialize chat history and conversation state
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    
    if "conversation_state" not in st.session_state:
        st.session_state.conversation_state = {
            "awaiting_filters": False,
            "awaiting_num_projects": False,
            "awaiting_extra_info": False,
            "person_name": None,
            "num_projects": None,
            "include_details": None,
            "all_projects": None,
            "filtered_projects": None
        }
    
    # Display example queries
    st.markdown("### 💡 Example Requests:")
    col1, col2 = st.columns(2)
    with col1:
        st.info("🔹 *Fill in the template with Bud van der Schier's most recent experience*")
        st.info("🔹 *Generate credentials for Tim Kramer*")
    with col2:
        st.info("🔹 *I need Charbel Moussa's latest project credentials*")
        st.info("🔹 *Create a presentation for Harold de Bruijn*")
    
    st.markdown("**💡 Filtering examples:**")
    st.caption("After selecting a person, you can filter by: year ranges, industries, sectors, service lines, or clients")
    
    st.markdown("---")
    
    # Chat interface
    st.subheader("💬 Chat with the Agent")
    
    # Display chat history
    for message in st.session_state.chat_history:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
            if "dataframe" in message and message["dataframe"] is not None:
                st.dataframe(message["dataframe"], width="stretch")
            if "download_data" in message:
                st.download_button(
                    label="📥 Download PowerPoint",
                    data=message["download_data"],
                    file_name=message["download_filename"],
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=message.get("download_key", "download")
                )
    
    # Chat input
    if prompt := st.chat_input("Type your request here..."):
        # Add user message
        st.session_state.chat_history.append({
            "role": "user",
            "content": prompt
        })
        
        with st.chat_message("user"):
            st.markdown(prompt)
        
        # Check if we're in a follow-up conversation state
        state = st.session_state.conversation_state
        
        # Process based on conversation state
        with st.chat_message("assistant"):
            
            # If waiting for filter criteria
            if state["awaiting_filters"]:
                user_filter = prompt.lower().strip()
                
                # Check if user wants to skip filtering
                if any(skip_word in user_filter for skip_word in ["no", "skip", "all", "none", "everything"]):
                    state["awaiting_filters"] = False
                    state["awaiting_num_projects"] = True
                    state["filtered_projects"] = state["all_projects"]
                    
                    response = f"""✅ Got it! Using all **{len(state['all_projects'])} projects**.

❓ How many projects would you like to include in the presentation?
   
💡 *Tip: For partners, 5-7 projects work well. For managers, 3-5 is typical.*

Please type a number (e.g., '5' or '7'):"""
                    
                    st.markdown(response)
                    st.session_state.chat_history.append({
                        "role": "assistant",
                        "content": response
                    })
                    st.rerun()
                
                else:
                    # Apply filters using LLM
                    with st.spinner("🔍 Applying filters..."):
                        filtered = apply_filters_with_llm(state["all_projects"], prompt)
                    
                    if len(filtered) == 0:
                        response = f"⚠️ No projects match those filters. Try different criteria or type 'skip' to use all {len(state['all_projects'])} projects."
                        st.markdown(response)
                        st.session_state.chat_history.append({
                            "role": "assistant",
                            "content": response
                        })
                        st.rerun()
                    else:
                        state["filtered_projects"] = filtered
                        state["awaiting_filters"] = False
                        state["awaiting_num_projects"] = True
                        
                        response = f"""✅ Filtered to **{len(filtered)} projects** (from {len(state['all_projects'])} total).

❓ How many projects would you like to include in the presentation?
   
💡 *Tip: For partners, 5-7 projects work well. For managers, 3-5 is typical.*

Please type a number (e.g., '5' or '7'):"""
                        
                        st.markdown(response)
                        st.session_state.chat_history.append({
                            "role": "assistant",
                            "content": response
                        })
                        st.rerun()
            
            # If waiting for number of projects
            if state["awaiting_num_projects"]:
                try:
                    # Try to extract number from user input
                    import re
                    numbers = re.findall(r'\d+', prompt)
                    if numbers:
                        num_projects = int(numbers[0])
                        state["num_projects"] = num_projects
                        state["awaiting_num_projects"] = False
                        state["awaiting_extra_info"] = True
                        
                        response = f"✅ Got it! I'll select **{num_projects} projects**.\n\n❓ Would you like me to include **detailed descriptions** for each project? (yes/no)"
                        st.markdown(response)
                        st.session_state.chat_history.append({
                            "role": "assistant",
                            "content": response
                        })
                    else:
                        response = "⚠️ I couldn't understand the number. Please type a number like '5' or '10'."
                        st.markdown(response)
                        st.session_state.chat_history.append({
                            "role": "assistant",
                            "content": response
                        })
                except Exception as e:
                    response = f"⚠️ Error: {str(e)}. Please type a number."
                    st.markdown(response)
                    st.session_state.chat_history.append({
                        "role": "assistant",
                        "content": response
                    })
                st.rerun()
            
            # If waiting for extra info confirmation
            elif state["awaiting_extra_info"]:
                user_response = prompt.lower().strip()
                
                if "yes" in user_response or "y" == user_response or "sure" in user_response or "ok" in user_response:
                    state["include_details"] = True
                    response_text = "✅ Perfect! I'll include detailed descriptions.\n\n🔄 Now generating your presentation..."
                elif "no" in user_response or "n" == user_response or "nope" in user_response:
                    state["include_details"] = False
                    response_text = "✅ Got it! I'll keep it concise without detailed descriptions.\n\n🔄 Now generating your presentation..."
                else:
                    response = "⚠️ Please answer with 'yes' or 'no'."
                    st.markdown(response)
                    st.session_state.chat_history.append({
                        "role": "assistant",
                        "content": response
                    })
                    st.rerun()
                
                st.markdown(response_text)
                
                # Now proceed with generation
                state["awaiting_extra_info"] = False
                person_name = state["person_name"]
                num_projects = state["num_projects"]
                include_details = state["include_details"]
                person_projects = state["filtered_projects"] if state["filtered_projects"] is not None else state["all_projects"]
                
                # Select best projects
                with st.spinner("🤖 Selecting best projects..."):
                    selected_projects, reasoning = select_best_projects(
                        person_projects, 
                        num_projects, 
                        "most recent"
                    )
                
                st.info(f"💡 **Selection reasoning**: {reasoning}")
                
                # Display selected projects
                st.markdown("### 📋 Selected Projects:")
                display_columns = [
                    "Project name / engagement title",
                    "Client name",
                    "Year of completion",
                    "Industry",
                    "Engagement partner",
                    "Engagement manager",
                    "Team members"
                ]
                # Only show columns that exist
                display_columns = [col for col in display_columns if col in selected_projects.columns]
                display_df = selected_projects[display_columns]
                st.dataframe(display_df, use_container_width=True)
                
                # Generate PowerPoint
                with st.spinner("📄 Generating PowerPoint..."):
                    # Try multiple paths for the template
                    possible_template_paths = [
                        'Bud van der Schier– Partner.pptx',  # Same directory as app.py
                        os.path.join(os.path.dirname(__file__), '..', 'Bud van der Schier– Partner.pptx'),  # Parent
                        os.path.join(os.path.dirname(os.path.dirname(__file__)), 'Bud van der Schier– Partner.pptx'),  # Two up
                        os.path.join(os.path.dirname(__file__), 'Bud van der Schier– Partner.pptx')  # Pages dir
                    ]
                    
                    template_path = None
                    for path in possible_template_paths:
                        if os.path.exists(path):
                            template_path = path
                            break
                    
                    if not template_path:
                        st.error("Template file 'Bud van der Schier– Partner.pptx' not found")
                        st.stop()
                    
                    try:
                        pptx_buffer = create_filled_presentation(
                            template_path,
                            selected_projects,
                            person_name,
                            include_details
                        )
                        
                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                        filename = f"Credentials_{person_name.replace(' ', '_')}_{timestamp}.pptx"
                        
                        st.success("✅ PowerPoint generated successfully!")
                        
                        download_key = f"download_{len(st.session_state.chat_history)}"
                        
                        st.download_button(
                            label="📥 Download PowerPoint",
                            data=pptx_buffer.getvalue(),
                            file_name=filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key=download_key
                        )
                        
                        # Store in chat history
                        full_response = f"""{response_text}

💡 **Selection reasoning**: {reasoning}

### 📋 Selected Projects:

✅ PowerPoint generated successfully!"""
                        
                        st.session_state.chat_history.append({
                            "role": "assistant",
                            "content": full_response,
                            "dataframe": display_df,
                            "download_data": pptx_buffer.getvalue(),
                            "download_filename": filename,
                            "download_key": download_key
                        })
                        
                        # Reset conversation state
                        st.session_state.conversation_state = {
                            "awaiting_filters": False,
                            "awaiting_num_projects": False,
                            "awaiting_extra_info": False,
                            "person_name": None,
                            "num_projects": None,
                            "include_details": None,
                            "all_projects": None,
                            "filtered_projects": None
                        }
                        
                    except Exception as e:
                        error_msg = f"❌ Error generating PowerPoint: {str(e)}"
                        st.error(error_msg)
                        st.session_state.chat_history.append({
                            "role": "assistant",
                            "content": error_msg
                        })
                
                st.rerun()
            
            # Initial request - parse and ask for details
            else:
                with st.spinner("🤔 Understanding your request..."):
                    parsed_request = parse_user_request(prompt, all_people)
                
                if parsed_request and parsed_request.get("found_match"):
                    person_name = parsed_request["person_name"]
                    
                    # Find projects
                    with st.spinner(f"🔍 Finding projects for {person_name}..."):
                        person_projects = find_person_projects(df, person_name)
                    
                    if len(person_projects) > 0:
                        # Store in conversation state
                        state["person_name"] = person_name
                        state["all_projects"] = person_projects
                        state["awaiting_filters"] = True
                        
                        response = f"""✅ Found **{len(person_projects)} projects** for **{person_name}**!

🔍 Would you like to **filter** these projects? 

**Examples:**
- *"Only projects from 2024-2025"*
- *"Only financial services projects"*
- *"Exclude audit projects"*
- *"Projects in healthcare or technology"*

Or type **'skip'** to use all projects:"""
                        
                        st.markdown(response)
                        st.session_state.chat_history.append({
                            "role": "assistant",
                            "content": response
                        })
                    else:
                        error_msg = f"❌ No projects found for {person_name}"
                        st.error(error_msg)
                        st.session_state.chat_history.append({
                            "role": "assistant",
                            "content": error_msg
                        })
                else:
                    error_msg = """❌ I couldn't understand your request.

Please try something like:
- "Fill in the template with Bud van der Schier's most recent experience"
- "Generate credentials for Tim Kramer"
- "Create presentation for Charbel Moussa"

Available people: """ + ", ".join(all_people[:10]) + "..."
                    
                    st.error(error_msg)
                    st.session_state.chat_history.append({
                        "role": "assistant",
                        "content": error_msg
                    })
                
                st.rerun()
    
    # Clear chat button
    if st.button("🗑️ Clear Chat History"):
        st.session_state.chat_history = []
        st.rerun()

else:
    st.error("❌ Failed to load credentials database")
    st.info("Please ensure 'credentials_full.xlsx' exists in the credentials folder")
